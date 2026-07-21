"""
Executive Summary slide factory for SQUAD presentations.
Generates the slide showing ROI totals, savings breakdown,
and a scaling tier table for different deployment amounts.
"""

from typing import Any, Dict

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig

# -- Layout constants --
HEADER_FONT_SIZE: int = 24
BODY_FONT_SIZE: int = 12
TABLE_FONT_SIZE: int = 11
BULLET_INDENT: float = 0.5

# -- Colors --
PURPLE_HEADER: RGBColor = RGBColor(75, 0, 130)
WHITE: RGBColor = RGBColor(255, 255, 255)
BLACK: RGBColor = RGBColor(0, 0, 0)

# -- Slide number label --
SLIDE_NUMBER: str = "04"

# -- Value proposition bullet points (static content) --
VALUE_PROPOSITIONS = [
    (
        "Real time capacity and risks assessment:",
        "Data-driven decision making in order to mitigate supply chain risks.",
    ),
    (
        "Asset & Life cycle management:",
        "Visibility into tools coming to their end of life or needing "
        "refurbishment & asset location and status globally.",
    ),
    (
        "Cycle time compliance:",
        "Track cycle time deviations and establish realistic demand planning "
        "based on actual cycle time.",
    ),
    (
        "Maintenance:",
        "Downtime reduction through maintenance schedules and optimization.",
    ),
    (
        "Cost reduction:",
        "Cycle time deviation from RFQ and PPAP, machine hour validation.",
    ),
    (
        "Capacity Management:",
        "Supply risk mitigation, utilization of real time Run-Rate analysis "
        "to track tooling efficiency, MTBF and MTTR.",
    ),
]


def add_executive_summary_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    totals: Dict[str, Any],
) -> None:
    """
    Add the Executive Summary slide to the presentation.

    Shows the purpose of engagement, value propositions, ROI bullet points,
    and a deployment scaling tier table.

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration.
        totals: Aggregated ROI totals from compute_executive_totals().
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    # Header bar
    _add_section_header(slide, SLIDE_NUMBER, "Executive Summary")

    # Purpose paragraph
    purpose_top = 1.3
    _add_purpose_block(slide, config, purpose_top)

    # Value propositions
    vp_top = 2.5
    _add_value_propositions(slide, vp_top)

    # ROI bullets (left column)
    roi_top = 5.0
    _add_roi_bullets(slide, totals, config, roi_top)

    # Tier table (right column)
    _add_tier_table(slide, totals, config, roi_top)


def _add_section_header(slide, number: str, title: str) -> None:
    """Add numbered section header bar at the top of the slide."""
    # Number badge
    badge = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.2), Inches(0.6), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0, 51, 102)
    badge.line.fill.background()

    badge_tf = badge.text_frame
    badge_tf.text = number
    bp = badge_tf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.font.size = Pt(16)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(0.2), Inches(7.0), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(HEADER_FONT_SIZE)
    tp.font.bold = True
    tp.font.color.rgb = BLACK


def _add_purpose_block(slide, config: SalesReportConfig, top: float) -> None:
    """Add the 'Purpose of Engagement' paragraph."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9.0), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"Purpose of Engagement with {config.client_name}:"
    p.font.size = Pt(14)
    p.font.bold = True

    desc = tf.add_paragraph()
    desc.text = (
        f"To address and solve {config.client_name}'s specific challenges in "
        "managing its extensive tooling supply chain by implementing a unified "
        "digital platform. This platform is designed to provide "
        f"{config.client_name} with real-time visibility over their tools, "
        "which is currently managed with static, and often unreliable, data."
    )
    desc.font.size = Pt(BODY_FONT_SIZE)
    desc.space_before = Pt(6)


def _add_value_propositions(slide, top: float) -> None:
    """Add the bullet list of solution value propositions."""
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(9.0), Inches(0.4)
    )
    hf = header_box.text_frame
    hf.text = "The eMoldino solution helps with:"
    hp = hf.paragraphs[0]
    hp.font.size = Pt(13)
    hp.font.bold = True

    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top + 0.4), Inches(9.0), Inches(2.2)
    )
    ctf = content_box.text_frame
    ctf.word_wrap = True

    for idx, (label, desc) in enumerate(VALUE_PROPOSITIONS):
        if idx == 0:
            p = ctf.paragraphs[0]
        else:
            p = ctf.add_paragraph()
        p.text = f"{label} {desc}"
        p.font.size = Pt(10)
        p.level = 1
        p.space_after = Pt(2)


def _add_roi_bullets(
    slide,
    totals: Dict[str, Any],
    config: SalesReportConfig,
    top: float,
) -> None:
    """Add the ROI breakdown bullets on the left side."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(4.0), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True

    header_p = tf.paragraphs[0]
    header_p.text = f"Benefits and returns experienced by {config.client_name}:"
    header_p.font.size = Pt(11)
    header_p.font.bold = True

    bullets = [
        ("Part Loss ROI", totals.get("part_loss_roi", 0)),
        ("Time Loss ROI", totals.get("time_loss_roi", 0)),
        ("CT Saving Improvement", totals.get("ct_roi", 0)),
        ("CT Opportunity", totals.get("ct_opportunity", 0)),
        ("Total Savings", totals.get("total_savings", 0)),
        ("Project Cost", totals.get("project_cost", 0)),
        ("ROI", f"~{totals.get('roi_ratio', 0):.2f}x"),
    ]

    for label, value in bullets:
        p = tf.add_paragraph()
        if isinstance(value, (int, float)):
            p.text = f"{label}: ${value:,.0f}"
        else:
            p.text = f"{label}: {value}"
        p.font.size = Pt(10)
        p.level = 1


def _add_tier_table(
    slide,
    totals: Dict[str, Any],
    config: SalesReportConfig,
    top: float,
) -> None:
    """Add the deployment scaling tier table on the right side."""
    tiers = list(config.roi_tiers.items())
    rows = len(tiers) + 1
    cols = 3

    table_shape = slide.shapes.add_table(
        rows, cols, Inches(5.0), Inches(top), Inches(4.5), Inches(1.5)
    )
    table = table_shape.table

    # Header row
    headers = [
        "Deployment Amount",
        "Expected Saving Opportunity",
        "Return on Investments (ROI)",
    ]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE_HEADER
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    base_savings = totals.get("total_savings", 0)
    tool_count = len(config.tools) if config.tools else 1

    for row_idx, (tier_count, tier_label) in enumerate(tiers, start=1):
        scale_factor = tier_count / max(tool_count, 1)
        projected_saving = base_savings * scale_factor

        table.cell(row_idx, 0).text = tier_label
        table.cell(row_idx, 1).text = f"${projected_saving:,.0f}"

        if row_idx == 1:
            table.cell(row_idx, 2).text = ""
        else:
            ratio = totals.get("roi_ratio", 0)
            table.cell(row_idx, 2).text = f"ROI: ~{ratio:.2f}x"

        for col_idx in range(cols):
            p = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
