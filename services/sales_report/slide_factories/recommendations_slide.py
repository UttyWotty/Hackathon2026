"""
Business Recommendations slide factory for SQUAD presentations.
Generates the 'What is needed from Client' slide with the ask amount,
supplier/tool scope, and expected savings (PDF page 5).
"""

from typing import Any, Dict

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import BusinessRecommendation, SalesReportConfig

# -- Layout constants --
HEADER_FONT_SIZE: int = 20
BODY_FONT_SIZE: int = 14

# -- Colors --
DARK_BLUE: RGBColor = RGBColor(0, 51, 102)
WHITE: RGBColor = RGBColor(255, 255, 255)
BLACK: RGBColor = RGBColor(0, 0, 0)
RED_BADGE: RGBColor = RGBColor(200, 16, 46)
LIGHT_BLUE_BG: RGBColor = RGBColor(189, 215, 238)


def add_recommendations_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    totals: Dict[str, Any],
) -> None:
    """
    Add the Business Recommendations slide.

    Shows the budget ask, supplier/tool scope, expected savings,
    and any custom notes from the sales team.

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration.
        totals: Executive totals for savings projections.
    """
    recs = config.recommendations or BusinessRecommendation()
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    _add_header(
        slide,
        "05",
        f"Business Recommendations - What is needed from {config.client_name}?",
    )

    # "The Ask" badge
    _add_ask_badge(slide)

    # Budget description
    _add_budget_block(slide, config, recs)

    # Action items
    _add_action_items(slide, config)

    # Savings box
    _add_savings_box(slide, recs, totals)


def _add_header(slide, number: str, title: str) -> None:
    """Add section header bar."""
    badge = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.2), Inches(0.6), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    bf = badge.text_frame
    bf.text = number
    bp = bf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(0.2), Inches(8.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(HEADER_FONT_SIZE)
    tp.font.bold = True


def _add_ask_badge(slide) -> None:
    """Add the red 'The Ask:' badge."""
    badge = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.0), Inches(1.2), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RED_BADGE
    badge.line.fill.background()
    bf = badge.text_frame
    bf.text = "The Ask:"
    bp = bf.paragraphs[0]
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = WHITE


def _add_budget_block(
    slide,
    config: SalesReportConfig,
    recs: BusinessRecommendation,
) -> None:
    """Add the budget description text."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = (
        f"Establish a budget to cover {recs.supplier_count} Critical "
        f"Suppliers and their existing tools."
    )
    p1.font.size = Pt(BODY_FONT_SIZE)
    p1.font.bold = True

    if recs.ask_amount > 0:
        p2 = tf.add_paragraph()
        p2.text = (
            f"Asking for ${recs.ask_amount:,.0f} to do "
            f"{recs.tool_count:,} tools across "
            f"{recs.supplier_count} suppliers "
            f"(based on {recs.tools_per_supplier} tool avg per site)."
        )
        p2.font.size = Pt(12)
        p2.level = 1
        p2.space_before = Pt(8)

    for note in recs.notes:
        pn = tf.add_paragraph()
        pn.text = note
        pn.font.size = Pt(12)
        pn.level = 1


def _add_action_items(slide, config: SalesReportConfig) -> None:
    """Add the standard action items list."""
    items = [
        "List of Critical Suppliers and Location in order to communicate "
        "the timeline to implement hardware and software.",
        "Establish a Project Management team.",
    ]
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9.0), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.level = 1


def _add_savings_box(
    slide,
    recs: BusinessRecommendation,
    totals: Dict[str, Any],
) -> None:
    """Add the expected savings summary box."""
    saving = recs.expected_saving or totals.get("total_savings", 0)

    box = slide.shapes.add_shape(1, Inches(3.0), Inches(5.0), Inches(4.5), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE_BG
    box.line.fill.background()

    text_box = slide.shapes.add_textbox(
        Inches(3.2), Inches(5.1), Inches(4.0), Inches(1.0)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "Expected Saving Opportunity"
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    if recs.ask_amount > 0:
        p2.text = f"The Ask: ${recs.ask_amount:,.0f}"
    else:
        p2.text = "The Ask: TBD"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = f"Expected Saving Opportunity: ${saving:,.0f}"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.alignment = PP_ALIGN.CENTER
    p3.font.color.rgb = RGBColor(0, 112, 192)
