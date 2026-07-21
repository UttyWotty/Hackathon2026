"""
Weekly Tooling Performance slide factory for SQUAD presentations.
Generates the week-over-week KPI comparison table and capacity risk
section for each tool (matching PDF pages 12, 14, 16, 18, 20).
"""

from typing import Any, Dict, List

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from analysis.shared.weekly_comparison_ppt import calculate_percentage_change
from services.sales_report.config import SalesReportConfig, ToolConfig

# -- Layout constants --
HEADER_FONT_SIZE: int = 18
TABLE_FONT_SIZE: int = 9

# -- Colors --
DARK_BLUE: RGBColor = RGBColor(0, 51, 102)
BLUE_HEADER: RGBColor = RGBColor(68, 114, 196)
WHITE: RGBColor = RGBColor(255, 255, 255)
BLACK: RGBColor = RGBColor(0, 0, 0)
GREEN: RGBColor = RGBColor(0, 128, 0)
RED: RGBColor = RGBColor(255, 0, 0)
LIGHT_GRAY_BG: RGBColor = RGBColor(220, 220, 220)

# -- KPI row definitions (label, key, unit) --
KPI_ROWS: List[tuple] = [
    ("Total Incurred Loss in Machine Hours:", "loss_hours", "hours"),
    ("Total Incurred Costs (Machine Rate + Labor)*", "incurred_costs", "dollars"),
    ("Parts Opportunity Lost", "parts_opportunity_lost", "parts"),
    ("Run Rate Efficiency", "efficiency_percentage", "percent"),
    ("Run Rate MTBF", "mtbf_minutes", "minutes"),
    ("Run Rate MTTR", "mttr_minutes", "minutes"),
]

CAPACITY_ROWS: List[tuple] = [
    ("Optimal Output (100% OEE)**", "optimal_output", "parts_decimal"),
    ("Actual Output", "actual_output", "parts_decimal"),
    ("Availability (Gain/Loss)", "availability_loss", "parts_signed"),
    ("Efficiency (Gain/Loss)", "efficiency_loss", "parts_signed"),
]

MAX_WEEKS: int = 4


def add_weekly_performance_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    tool_cfg: ToolConfig,
    weekly_kpis: List[Dict[str, Any]],
    month_label: str = "",
) -> None:
    """
    Add a weekly tooling performance comparison slide for one tool.

    Displays a table with up to 4 weeks of KPI data plus a monthly total
    column, followed by a capacity risk section.

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration.
        tool_cfg: Tool configuration.
        weekly_kpis: List of KPI dicts, one per week (up to 4).
        month_label: Label for the monthly total column header.
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    title = f"Tooling Performance: {tool_cfg.equipment_code} - " f"{tool_cfg.commodity}"
    _add_title_block(slide, title)

    num_weeks = min(len(weekly_kpis), MAX_WEEKS)
    _add_kpi_table(slide, weekly_kpis[:num_weeks], month_label)
    _add_capacity_section(slide, weekly_kpis[:num_weeks], month_label)
    _add_footnotes(slide)


def _add_title_block(slide, title: str) -> None:
    """Add the VANTIS-branded title block."""
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.2), Inches(8.0), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(HEADER_FONT_SIZE)
    p.font.bold = True
    p.font.color.rgb = BLACK

    subtitle_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.7), Inches(8.0), Inches(0.3)
    )
    sf = subtitle_box.text_frame
    sf.text = "{Comparison between previous week}"
    sp = sf.paragraphs[0]
    sp.font.size = Pt(12)
    sp.font.italic = True


def _add_kpi_table(
    slide,
    weekly_kpis: List[Dict[str, Any]],
    month_label: str,
) -> None:
    """Build the main KPI comparison table."""
    num_weeks = len(weekly_kpis)
    cols = num_weeks + 2  # KPI label + weeks + monthly total
    rows = len(KPI_ROWS) + 1  # header + data

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.3),
        Inches(1.2),
        Inches(9.4),
        Inches(2.8),
    )
    table = table_shape.table

    # Header
    table.cell(0, 0).text = "KPI"
    for w in range(num_weeks):
        table.cell(0, w + 1).text = f"Week {w + 1}"
    monthly_header = (
        f"Monthly Total\n({month_label})" if month_label else "Monthly Total"
    )
    table.cell(0, cols - 1).text = monthly_header

    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_HEADER
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(TABLE_FONT_SIZE)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r_idx, (label, key, unit) in enumerate(KPI_ROWS, start=1):
        table.cell(r_idx, 0).text = label
        p = table.cell(r_idx, 0).text_frame.paragraphs[0]
        p.font.size = Pt(TABLE_FONT_SIZE)

        for w in range(num_weeks):
            kpi = weekly_kpis[w]
            value = kpi.get(key, "No Production")
            cell_text = _format_value(value, unit)

            # Add percentage change vs previous week
            if w > 0 and isinstance(value, (int, float)):
                prev = weekly_kpis[w - 1].get(key)
                if isinstance(prev, (int, float)) and prev != 0:
                    pct = calculate_percentage_change(value, prev)
                    if pct:
                        cell_text = f"{cell_text} ({pct})"

            cell = table.cell(r_idx, w + 1)
            cell.text = cell_text
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(TABLE_FONT_SIZE)
            cp.alignment = PP_ALIGN.CENTER

        # Monthly total (average for rates, sum for quantities)
        _set_monthly_total(table, r_idx, cols - 1, weekly_kpis, key, unit)


def _add_capacity_section(
    slide,
    weekly_kpis: List[Dict[str, Any]],
    month_label: str,
) -> None:
    """Add the capacity risk sub-table below the KPI table."""
    num_weeks = len(weekly_kpis)
    cols = num_weeks + 2
    rows = len(CAPACITY_ROWS) + 1

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.3),
        Inches(4.2),
        Inches(9.4),
        Inches(1.8),
    )
    table = table_shape.table

    # Header
    header_cell = table.cell(0, 0)
    header_cell.merge(table.cell(0, cols - 1))
    header_cell.text = "Capacity Risk"
    header_cell.fill.solid()
    header_cell.fill.fore_color.rgb = DARK_BLUE
    hp = header_cell.text_frame.paragraphs[0]
    hp.font.size = Pt(TABLE_FONT_SIZE)
    hp.font.bold = True
    hp.font.color.rgb = WHITE

    # Data rows (capacity values are typically monthly, span all columns)
    for r_idx, (label, key, unit) in enumerate(CAPACITY_ROWS, start=1):
        table.cell(r_idx, 0).text = label
        p = table.cell(r_idx, 0).text_frame.paragraphs[0]
        p.font.size = Pt(TABLE_FONT_SIZE)

        # Merge data columns for capacity (single value spans the row)
        merged_cell = table.cell(r_idx, 1)
        merged_cell.merge(table.cell(r_idx, cols - 1))

        combined = _get_monthly_capacity_value(weekly_kpis, key, unit)
        merged_cell.text = combined
        mp = merged_cell.text_frame.paragraphs[0]
        mp.font.size = Pt(TABLE_FONT_SIZE)
        mp.alignment = PP_ALIGN.CENTER
        mp.font.bold = True


def _add_footnotes(slide) -> None:
    """Add footnote text at the bottom."""
    box = slide.shapes.add_textbox(Inches(0.3), Inches(6.3), Inches(9.0), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    notes = [
        "*based on a machine rate cost of $170/h and labor costs of $10/h",
        "** Target OEE score to be input by client for future reports.",
        "*** Average of all weeks",
    ]
    for idx, note in enumerate(notes):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = note
        p.font.size = Pt(7)
        p.font.italic = True


def _format_value(value: Any, unit: str) -> str:
    """Format a KPI value based on its unit type."""
    if value is None or value == "No Production":
        return "No Production"
    if unit == "hours":
        hours = float(value)
        h = int(hours)
        m = int((hours - h) * 60)
        if h > 0:
            return f"{h} Hour{'s' if h != 1 else ''}" + (f" {m}M" if m else "")
        return f"{m}M"
    if unit == "dollars":
        return f"${float(value):,.0f}"
    if unit == "parts":
        return f"{int(value):,} Parts"
    if unit == "percent":
        return f"{float(value):.1f}%"
    if unit == "minutes":
        return f"{float(value):.1f} Minutes"
    return str(value)


def _set_monthly_total(
    table,
    row_idx: int,
    col_idx: int,
    weekly_kpis: List[Dict[str, Any]],
    key: str,
    unit: str,
) -> None:
    """Compute and set the monthly total/average cell."""
    values = [w.get(key) for w in weekly_kpis if isinstance(w.get(key), (int, float))]
    if not values:
        table.cell(row_idx, col_idx).text = "N/A"
    elif unit in ("percent", "minutes"):
        avg = sum(values) / len(values)
        label = _format_value(avg, unit)
        table.cell(row_idx, col_idx).text = f"{label}***"
    else:
        total = sum(values)
        table.cell(row_idx, col_idx).text = _format_value(total, unit)

    p = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
    p.font.size = Pt(TABLE_FONT_SIZE)
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True


def _get_monthly_capacity_value(
    weekly_kpis: List[Dict[str, Any]],
    key: str,
    unit: str,
) -> str:
    """Extract or compute a capacity value across weeks."""
    values = [w.get(key) for w in weekly_kpis if isinstance(w.get(key), (int, float))]
    if not values:
        return "N/A"
    total = sum(values)
    if unit == "parts_decimal":
        return f"{total:,.3f} Parts"
    if unit == "parts_signed":
        sign = "+" if total > 0 else ""
        return f"{sign}{total:,.0f} Parts"
    return str(total)
