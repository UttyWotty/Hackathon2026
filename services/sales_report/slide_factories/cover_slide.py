"""
Cover slide factory for SQUAD presentations.
Generates the opening slide with client branding, presentation title,
date, and eMoldino/client logos.
"""

import os

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig

# -- Layout constants --
SLIDE_WIDTH_INCHES: float = 10.0
TITLE_LEFT: float = 0.5
TITLE_TOP: float = 1.5
TITLE_WIDTH: float = 5.0
TITLE_HEIGHT: float = 1.2
SUBTITLE_TOP: float = 2.8
SUBTITLE_HEIGHT: float = 0.5
LOGO_SIZE: float = 0.8
LOGO_MARGIN: float = 0.3
TITLE_FONT_SIZE: int = 36
SUBTITLE_FONT_SIZE: int = 18
WHITE: RGBColor = RGBColor(255, 255, 255)


def add_cover_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    report_date_label: str = "",
) -> None:
    """
    Add the cover slide to the presentation.

    Places the presentation title and date on the left, with optional
    background image and logos in the top-right corner.

    Args:
        ppt: PPTGenerator instance to add the slide to.
        config: Sales report configuration with client info and assets.
        report_date_label: Formatted date string (e.g. "March, 2026").
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    # Background image (full slide)
    cover_path = os.path.join(config.assets_base_dir, config.assets.cover_image)
    if config.assets.cover_image and os.path.exists(cover_path):
        slide.shapes.add_picture(
            cover_path,
            Inches(0),
            Inches(0),
            width=ppt.prs.slide_width,
            height=ppt.prs.slide_height,
        )

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(TITLE_LEFT),
        Inches(TITLE_TOP),
        Inches(TITLE_WIDTH),
        Inches(TITLE_HEIGHT),
    )
    tf = title_box.text_frame
    tf.text = config.presentation_title
    p = tf.paragraphs[0]
    p.font.size = Pt(TITLE_FONT_SIZE)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Date subtitle
    if report_date_label:
        sub_box = slide.shapes.add_textbox(
            Inches(TITLE_LEFT),
            Inches(SUBTITLE_TOP),
            Inches(TITLE_WIDTH),
            Inches(SUBTITLE_HEIGHT),
        )
        sf = sub_box.text_frame
        sf.text = report_date_label
        sp = sf.paragraphs[0]
        sp.font.size = Pt(SUBTITLE_FONT_SIZE)
        sp.font.color.rgb = WHITE

    # Logos (top-right)
    _add_logos(slide, config)


def _add_logos(slide, config: SalesReportConfig) -> None:
    """
    Place eMoldino and client logos in the top-right corner.

    Args:
        slide: PowerPoint slide object.
        config: Sales report configuration with asset paths.
    """
    x_pos = SLIDE_WIDTH_INCHES - LOGO_MARGIN - LOGO_SIZE

    emoldino_path = os.path.join(config.assets_base_dir, config.assets.emoldino_logo)
    if os.path.exists(emoldino_path):
        slide.shapes.add_picture(
            emoldino_path,
            Inches(x_pos - LOGO_SIZE - LOGO_MARGIN),
            Inches(LOGO_MARGIN),
            width=Inches(LOGO_SIZE),
            height=Inches(LOGO_SIZE),
        )

    client_logo_path = os.path.join(config.assets_base_dir, config.assets.client_logo)
    if config.assets.client_logo and os.path.exists(client_logo_path):
        slide.shapes.add_picture(
            client_logo_path,
            Inches(x_pos),
            Inches(LOGO_MARGIN),
            width=Inches(LOGO_SIZE),
            height=Inches(LOGO_SIZE),
        )
