"""
Table of Contents slide factory for SQUAD presentations.
Generates a TOC slide with section labels matching the presentation structure.
"""

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig

# -- Layout constants --
SECTION_BOX_WIDTH: float = 5.0
SECTION_BOX_HEIGHT: float = 0.6
SECTION_X_START: float = 4.5
SECTION_Y_START: float = 2.0
SECTION_Y_GAP: float = 0.8
TITLE_FONT_SIZE: int = 28
SECTION_FONT_SIZE: int = 16

# -- Section colors --
COLOR_DARK_GRAY: RGBColor = RGBColor(89, 89, 89)
COLOR_LIGHT_BLUE: RGBColor = RGBColor(189, 215, 238)
WHITE: RGBColor = RGBColor(255, 255, 255)

# -- Default sections --
DEFAULT_SECTIONS = [
    ("eMoldino Corporate Overview", COLOR_DARK_GRAY, WHITE),
    ("Executive Summary", COLOR_LIGHT_BLUE, RGBColor(0, 0, 0)),
    ("Financial Results", COLOR_LIGHT_BLUE, RGBColor(0, 0, 0)),
    ("Business Recommendations", COLOR_LIGHT_BLUE, RGBColor(0, 0, 0)),
]


def add_toc_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
) -> None:
    """
    Add the Table of Contents slide to the presentation.

    Renders section labels as colored rectangles aligned on the right side,
    with optional background image on the left.

    Args:
        ppt: PPTGenerator instance to add the slide to.
        config: Sales report configuration (for background assets).
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    # Slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(4.0), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.text = "Table of Contents"
    p = tf.paragraphs[0]
    p.font.size = Pt(TITLE_FONT_SIZE)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Section boxes
    for idx, (label, bg_color, text_color) in enumerate(DEFAULT_SECTIONS):
        y = SECTION_Y_START + idx * SECTION_Y_GAP
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(SECTION_X_START),
            Inches(y),
            Inches(SECTION_BOX_WIDTH),
            Inches(SECTION_BOX_HEIGHT),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

        text_box = slide.shapes.add_textbox(
            Inches(SECTION_X_START + 0.2),
            Inches(y + 0.1),
            Inches(SECTION_BOX_WIDTH - 0.4),
            Inches(SECTION_BOX_HEIGHT - 0.2),
        )
        stf = text_box.text_frame
        stf.text = label
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(SECTION_FONT_SIZE)
        sp.font.bold = True
        sp.font.color.rgb = text_color

    # Footer
    _add_footer(slide, config)


def _add_footer(slide, config: SalesReportConfig) -> None:
    """
    Add confidential footer and logos to the bottom of the slide.

    Args:
        slide: PowerPoint slide object.
        config: Sales report configuration.
    """
    footer_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.0), Inches(6.0), Inches(0.3)
    )
    ff = footer_box.text_frame
    ff.text = (
        "Confidential and Proprietary. Copyright 2026. eMoldino. "
        "All rights reserved."
    )
    fp = ff.paragraphs[0]
    fp.font.size = Pt(8)
    fp.font.color.rgb = RGBColor(150, 150, 150)
