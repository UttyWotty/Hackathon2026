"""
Run Rate Efficiency table slide factory for SQUAD presentations.
Generates per-tool slides with tool metadata, ROI summary boxes,
and a monthly breakdown table of part loss and time loss (PDF pages 8-11).
"""

from typing import Any, Dict, List

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig, ToolConfig

# -- Layout constants --
HEADER_FONT_SIZE: int = 20
TABLE_FONT_SIZE: int = 9
META_FONT_SIZE: int = 10

# -- Colors --
DARK_BLUE: RGBColor = RGBColor(0, 51, 102)
GOLD_BG: RGBColor = RGBColor(255, 242, 204)
WHITE: RGBColor = RGBColor(255, 255, 255)
BLACK: RGBColor = RGBColor(0, 0, 0)
BLUE_HEADER: RGBColor = RGBColor(68, 114, 196)

# -- Monthly table row labels --
MONTHLY_ROW_LABELS: List[str] = [
    "Part Loss (Quantity)",
    "Time Loss (Downtime)",
    "Time Loss in Hours",
    "Time Loss in Seconds",
]


def add_run_rate_efficiency_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    tool_cfg: ToolConfig,
    savings: Dict[str, Any],
    slide_number: str = "08",
) -> None:
    """
    Add a Run Rate Efficiency table slide for a single tool.

    Shows tool metadata box, ROI summary cards, and monthly breakdown table.

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration.
        tool_cfg: Configuration for the specific tool.
        savings: Pre-computed savings dict from compute_tool_savings().
        slide_number: Slide number label.
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    title = (
        f"Run Deep Dive - Tool: {tool_cfg.equipment_code} " f"({tool_cfg.commodity})"
    )
    _add_header(slide, slide_number, title)

    metrics = savings

    # Tool metadata box (top-left)
    _add_tool_metadata_box(slide, tool_cfg)

    # ROI summary cards (top-right)
    _add_roi_summary_cards(slide, metrics, config)

    # Monthly breakdown table
    _add_monthly_table(slide, metrics)

    # Cost footnote
    _add_cost_footnote(slide, config, tool_cfg)


def _add_header(slide, number: str, title: str) -> None:
    """Add numbered header bar."""
    badge = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.2), Inches(0.6), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    bf = badge.text_frame
    bf.text = number
    bp = bf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(0.2), Inches(8.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(HEADER_FONT_SIZE)
    tp.font.bold = True


def _add_tool_metadata_box(slide, tool_cfg: ToolConfig) -> None:
    """Add the tool info box (Tool, Cavities, Commodity, CT)."""
    rows = 5
    cols = 2
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.2), Inches(3.0), Inches(1.8)
    )
    table = table_shape.table

    # Title row
    cell = table.cell(0, 0)
    cell.merge(table.cell(0, 1))
    cell.text = "Run Rate Efficiency"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(META_FONT_SIZE)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    meta_rows = [
        ("Tool", tool_cfg.equipment_code),
        ("Cavities", str(tool_cfg.cavities or "N/A")),
        ("Commodity", tool_cfg.commodity),
        (
            "Contracted Cycle Time (Seconds)",
            str(tool_cfg.contracted_ct_seconds or "N/A"),
        ),
    ]

    for idx, (label, value) in enumerate(meta_rows, start=1):
        table.cell(idx, 0).text = label
        table.cell(idx, 1).text = value
        for col in range(2):
            p = table.cell(idx, col).text_frame.paragraphs[0]
            p.font.size = Pt(META_FONT_SIZE)


def _add_roi_summary_cards(
    slide,
    metrics: Dict[str, Any],
    config: SalesReportConfig,
) -> None:
    """Add the ROI summary cards (Total Savings, Project Cost, ROI Ratio, etc.)."""
    total_savings = metrics.get("total_savings", 0)
    project_cost = config.project_cost
    roi_ratio = total_savings / project_cost if project_cost > 0 else 0.0

    left_items = [
        ("Total Savings", f"${total_savings:,.0f}"),
        ("Project Cost", f"${project_cost:,.0f}"),
        ("ROI Ratio", f"{roi_ratio:.2f}"),
    ]

    right_items = [
        ("Part Loss ROI", f"${metrics.get('part_loss_roi', 0):,.0f}"),
        ("Time Loss ROI", f"${metrics.get('time_loss_roi', 0):,.0f}"),
        ("CT ROI", f"${metrics.get('ct_opportunity', 0):,.0f}"),
    ]

    # Left card
    _draw_mini_table(slide, left_items, Inches(4.0), Inches(1.2))
    # Right card
    _draw_mini_table(slide, right_items, Inches(7.0), Inches(1.2))


def _draw_mini_table(slide, items, left, top) -> None:
    """Draw a small label-value table."""
    rows = len(items)
    cols = 2
    table_shape = slide.shapes.add_table(
        rows, cols, left, top, Inches(2.5), Inches(0.3 * rows)
    )
    table = table_shape.table

    for idx, (label, value) in enumerate(items):
        table.cell(idx, 0).text = label
        table.cell(idx, 1).text = value
        cell_bg = GOLD_BG
        table.cell(idx, 0).fill.solid()
        table.cell(idx, 0).fill.fore_color.rgb = cell_bg
        table.cell(idx, 1).fill.solid()
        table.cell(idx, 1).fill.fore_color.rgb = cell_bg
        for col in range(2):
            p = table.cell(idx, col).text_frame.paragraphs[0]
            p.font.size = Pt(META_FONT_SIZE)
            if col == 1:
                p.font.bold = True


def _add_monthly_table(slide, metrics: Dict[str, Any]) -> None:
    """Add the monthly breakdown table."""
    monthly_data = metrics.get("monthly_breakdown", {})
    months = sorted(monthly_data.keys()) if monthly_data else []

    cols = len(months) + 2  # label col + months + total col
    rows = len(MONTHLY_ROW_LABELS) + 1  # header + data rows

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(3.5),
        Inches(9.0),
        Inches(2.0),
    )
    table = table_shape.table

    # Header row
    table.cell(0, 0).text = "Month"
    for col_idx, month in enumerate(months, start=1):
        table.cell(0, col_idx).text = month
    table.cell(0, cols - 1).text = "Total Gain"

    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_HEADER
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(TABLE_FONT_SIZE)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, label in enumerate(MONTHLY_ROW_LABELS, start=1):
        table.cell(row_idx, 0).text = label
        total = 0
        for col_idx, month in enumerate(months, start=1):
            month_data = monthly_data.get(month, {})
            value = month_data.get(label.lower().replace(" ", "_"), 0)
            table.cell(row_idx, col_idx).text = str(value)
            if isinstance(value, (int, float)):
                total += value
        table.cell(row_idx, cols - 1).text = str(total)

        for col_idx in range(cols):
            p = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
            p.font.size = Pt(TABLE_FONT_SIZE)
            p.alignment = PP_ALIGN.CENTER


def _add_cost_footnote(
    slide,
    config: SalesReportConfig,
    tool_cfg: ToolConfig,
) -> None:
    """Add the industry-average cost footnote at the bottom."""
    part_cost = tool_cfg.resolved_part_cost()
    footnote_text = (
        f"Based on Industry Average\n"
        f"Machine Rate (Hourly): ${config.machine_rate_per_hour:,.0f}\n"
        f"Labor Cost: ${config.labor_cost_per_hour:,.0f}\n"
        f"Part Cost {tool_cfg.commodity}: ${part_cost}"
    )
    box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(3.5), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = footnote_text
    for p in tf.paragraphs:
        p.font.size = Pt(8)
        p.font.italic = True
