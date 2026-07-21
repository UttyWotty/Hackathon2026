"""
Run Deep Dive slide factory for SQUAD presentations.
Generates per-tool slides showing a screenshot placeholder on the left
and capacity-based savings calculations on the right (matching PDF page 7).
"""

import os
from typing import Any, Dict

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig, ToolConfig

# -- Layout constants --
HEADER_FONT_SIZE: int = 20
METRIC_FONT_SIZE: int = 12
PLACEHOLDER_TEXT: str = "[Insert platform screenshot here]"

# -- Colors --
WHITE: RGBColor = RGBColor(255, 255, 255)
BLACK: RGBColor = RGBColor(0, 0, 0)
DARK_BLUE: RGBColor = RGBColor(0, 51, 102)
LIGHT_BLUE_BG: RGBColor = RGBColor(240, 248, 255)
GRAY: RGBColor = RGBColor(150, 150, 150)


def add_run_deep_dive_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    tool_cfg: ToolConfig,
    savings: Dict[str, Any],
    slide_number: str = "07",
) -> None:
    """
    Add a Run Deep Dive slide for a single tool.

    Left side: screenshot placeholder or actual image.
    Right side: production time, downtime, part loss, savings metrics.

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration.
        tool_cfg: Configuration for the specific tool.
        savings: Pre-computed savings dict from compute_tool_savings().
        slide_number: Slide number label.
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    commodity = tool_cfg.commodity
    title_text = f"Run Deep Dive - Tool: {tool_cfg.equipment_code} ({commodity})"
    _add_header(slide, slide_number, title_text)

    # Left: screenshot or placeholder
    screenshot_key = f"deep_dive_{tool_cfg.equipment_code}"
    screenshot_path = config.screenshot_paths.get(screenshot_key, "")
    _add_screenshot_area(slide, screenshot_path)

    # Right: metrics from pre-computed savings
    _add_metrics_panel(slide, savings, tool_cfg, config)


def _add_header(slide, number: str, title: str) -> None:
    """Add the numbered header bar."""
    badge = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.2), Inches(0.6), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    bf = badge.text_frame
    bf.text = number
    bp = bf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(0.2), Inches(8.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(HEADER_FONT_SIZE)
    tp.font.bold = True


def _add_screenshot_area(slide, screenshot_path: str) -> None:
    """Add screenshot image or gray placeholder box."""
    left = Inches(0.3)
    top = Inches(1.0)
    width = Inches(5.5)
    height = Inches(5.5)

    if screenshot_path and os.path.exists(screenshot_path):
        slide.shapes.add_picture(screenshot_path, left, top, width=width, height=height)
    else:
        placeholder = slide.shapes.add_shape(1, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        placeholder.line.color.rgb = GRAY

        text_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.5), Inches(3.0), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.text = PLACEHOLDER_TEXT
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.font.italic = True


def _add_metrics_panel(
    slide,
    metrics: Dict[str, Any],
    tool_cfg: ToolConfig,
    config: SalesReportConfig,
) -> None:
    """Add the right-side panel with capacity/savings metrics."""
    right_x = 6.2
    top = 1.0

    # Cost assumption box
    part_cost = tool_cfg.resolved_part_cost()
    assumption_box = slide.shapes.add_textbox(
        Inches(right_x), Inches(top), Inches(3.5), Inches(0.4)
    )
    af = assumption_box.text_frame
    af.text = f"Assuming ${part_cost:.2f} for each {tool_cfg.commodity} Part:"
    ap = af.paragraphs[0]
    ap.font.size = Pt(10)
    ap.font.bold = True

    # Production / Downtime summary (from pre-computed savings)
    prod_minutes = metrics.get("production_minutes", 0)
    down_minutes = metrics.get("downtime_minutes", 0)
    prod_h, prod_m = divmod(int(prod_minutes), 60)
    down_h, down_m = divmod(int(down_minutes), 60)

    summary_items = [
        f"Total Production Time: {prod_h} Hours {prod_m} Minutes",
        f"Total Downtime: {down_h} Hours {down_m} Minutes",
    ]

    summary_box = slide.shapes.add_textbox(
        Inches(right_x), Inches(top + 0.6), Inches(3.5), Inches(1.0)
    )
    stf = summary_box.text_frame
    stf.word_wrap = True
    for idx, item in enumerate(summary_items):
        if idx == 0:
            p = stf.paragraphs[0]
        else:
            p = stf.add_paragraph()
        p.text = item
        p.font.size = Pt(METRIC_FONT_SIZE)
        p.font.bold = True

    # Savings bullets (already computed in dollar values)
    parts_lost = metrics.get("parts_lost", 0)
    part_loss_dollars = metrics.get("part_loss_roi", 0)
    time_loss_dollars = metrics.get("time_loss_roi", 0)
    ct_opportunity = metrics.get("ct_opportunity", 0)
    total_savings = metrics.get("total_savings", 0)

    savings_items = [
        f"Parts Lost Quantity: {int(parts_lost):,}",
        f"Part Loss in Dollars: ${part_loss_dollars:,.0f}",
        f"Time Loss in Dollars: ${time_loss_dollars:,.0f}",
        f"CT Opportunity: ${ct_opportunity:,.0f}",
        f"Total Savings: ${total_savings:,.0f}",
    ]

    savings_box = slide.shapes.add_textbox(
        Inches(right_x), Inches(top + 2.0), Inches(3.5), Inches(2.5)
    )
    svtf = savings_box.text_frame
    svtf.word_wrap = True
    for idx, item in enumerate(savings_items):
        if idx == 0:
            p = svtf.paragraphs[0]
        else:
            p = svtf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.level = 1
