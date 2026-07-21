"""
Appendix divider slide factory for SQUAD presentations.
Generates a simple dark-background slide with 'Appendix' title,
serving as a section break before detailed per-tool slides.
"""

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig

# -- Colors --
DARK_BG: RGBColor = RGBColor(10, 20, 40)
WHITE: RGBColor = RGBColor(255, 255, 255)
BLUE_LINE: RGBColor = RGBColor(0, 112, 192)

TITLE_FONT_SIZE: int = 44


def add_appendix_divider_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
) -> None:
    """
    Add an Appendix divider slide to the presentation.

    Dark background with large 'Appendix' title and a blue underline,
    plus footer logos.

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration (for logo assets).
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    # Dark background
    bg = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        ppt.prs.slide_width,
        ppt.prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5), Inches(5.0), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.text = "Appendix"
    p = tf.paragraphs[0]
    p.font.size = Pt(TITLE_FONT_SIZE)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Blue underline
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(3.8), Inches(3.0), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE_LINE
    line.line.fill.background()

    # Footer with confidential text
    footer_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(6.8), Inches(6.0), Inches(0.3)
    )
    ff = footer_box.text_frame
    ff.text = (
        "Confidential and Proprietary. Copyright 2026. eMoldino. All rights reserved."
    )
    fp = ff.paragraphs[0]
    fp.font.size = Pt(8)
    fp.font.color.rgb = RGBColor(150, 150, 150)
