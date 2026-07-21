"""
Run Rate Cycle Time chart slide factory for SQUAD presentations.
Generates per-tool slides with run rate efficiency, stability index,
and a screenshot placeholder for the cycle time chart (PDF pages 13, 15, 17, 19).
"""

import os
from typing import Any, Dict

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from analysis.shared.ppt_generator import PPTGenerator
from services.sales_report.config import SalesReportConfig, ToolConfig

# -- Layout constants --
HEADER_FONT_SIZE: int = 18
KPI_BOX_WIDTH: float = 3.0
KPI_BOX_HEIGHT: float = 0.4

# -- Colors --
DARK_BLUE: RGBColor = RGBColor(0, 51, 102)
WHITE: RGBColor = RGBColor(255, 255, 255)
BLACK: RGBColor = RGBColor(0, 0, 0)
LIGHT_GRAY: RGBColor = RGBColor(230, 230, 230)
GRAY_TEXT: RGBColor = RGBColor(150, 150, 150)
GOLD_BG: RGBColor = RGBColor(255, 242, 204)

PLACEHOLDER_TEXT: str = "[Insert Run Rate Cycle Time chart here]"


def add_run_rate_chart_slide(
    ppt: PPTGenerator,
    config: SalesReportConfig,
    tool_cfg: ToolConfig,
    runrate_result: Dict[str, Any],
    date_label: str = "",
    slide_number: str = "10",
) -> None:
    """
    Add a Run Rate Cycle Time chart slide for a single tool.

    Shows efficiency and stability KPI boxes at top, with a large
    chart area below (screenshot or placeholder).

    Args:
        ppt: PPTGenerator instance.
        config: Sales report configuration.
        tool_cfg: Tool configuration.
        runrate_result: RunRate analysis result for this tool.
        date_label: Date label (e.g. "November 14th, 2025").
        slide_number: Slide number label.
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    title = f"Tooling Performance: {tool_cfg.equipment_code} - " f"{tool_cfg.commodity}"
    _add_title(slide, slide_number, title)

    metrics = runrate_result.get("metrics", {})

    # Date badge
    if date_label:
        _add_date_badge(slide, date_label)

    # KPI boxes (RunRate tool returns "efficiency_percentage")
    efficiency = metrics.get("efficiency_percentage", 0)
    stability = metrics.get("stability_index", 0)
    _add_kpi_box(slide, 0.5, 1.5, f"Run Rate Efficiency: {efficiency:.1f}%")
    _add_kpi_box(slide, 6.0, 1.5, f"Run Rate Stability Index: {stability:.1f}%")

    # Descriptions
    _add_kpi_description(
        slide,
        0.5,
        2.1,
        "Indicating the percentage of shots that fell into the mode "
        'tolerance ("normal" shots)',
    )
    _add_kpi_description(
        slide,
        6.0,
        2.1,
        "Indicating the percentage of the total run that was spent in "
        '"normal" production',
    )

    # Chart area
    screenshot_key = f"runrate_chart_{tool_cfg.equipment_code}"
    screenshot_path = config.screenshot_paths.get(screenshot_key, "")
    _add_chart_area(slide, screenshot_path)


def _add_title(slide, number: str, title: str) -> None:
    """Add slide header."""
    badge = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.2), Inches(0.6), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    bf = badge.text_frame
    bf.text = number
    bp = bf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(0.2), Inches(8.5), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    tp = tf.paragraphs[0]
    tp.font.size = Pt(HEADER_FONT_SIZE)
    tp.font.bold = True


def _add_date_badge(slide, label: str) -> None:
    """Add centered date badge."""
    badge = slide.shapes.add_shape(
        1, Inches(3.5), Inches(1.0), Inches(3.0), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD_BG
    badge.line.fill.background()

    text_box = slide.shapes.add_textbox(
        Inches(3.5), Inches(1.0), Inches(3.0), Inches(0.35)
    )
    tf = text_box.text_frame
    tf.text = f"Date: {label}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True


def _add_kpi_box(slide, x: float, y: float, text: str) -> None:
    """Add a KPI summary box."""
    box = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(KPI_BOX_WIDTH), Inches(KPI_BOX_HEIGHT)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = BLACK

    text_box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(KPI_BOX_WIDTH), Inches(KPI_BOX_HEIGHT)
    )
    tf = text_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True


def _add_kpi_description(slide, x: float, y: float, text: str) -> None:
    """Add description text below a KPI box."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.5), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)


def _add_chart_area(slide, screenshot_path: str) -> None:
    """Add the chart screenshot or placeholder."""
    left = Inches(0.5)
    top = Inches(3.0)
    width = Inches(9.0)
    height = Inches(4.0)

    if screenshot_path and os.path.exists(screenshot_path):
        slide.shapes.add_picture(screenshot_path, left, top, width=width, height=height)
    else:
        placeholder = slide.shapes.add_shape(1, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = LIGHT_GRAY
        placeholder.line.color.rgb = GRAY_TEXT

        text_box = slide.shapes.add_textbox(
            Inches(3.0), Inches(4.8), Inches(4.0), Inches(0.5)
        )
        tf = text_box.text_frame
        tf.text = PLACEHOLDER_TEXT
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_TEXT
        p.font.italic = True
