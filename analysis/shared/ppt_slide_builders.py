"""PowerPoint slide building helpers for weekly comparison presentations.

This module contains all table construction, formatting, and slide layout
functions used by the weekly comparison PPT generator. It handles KPI
comparison tables, capacity detail slides, and cell-level formatting.
"""

from typing import Any, Dict

from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.util import Inches, Pt  # type: ignore

from analysis.shared.ppt_generator import PPTGenerator

# Week label constants (shared with main module)
LABEL_WEEK_1 = "Week 1"
LABEL_WEEK_2 = "Week 2"

# Table dimension constants
KPI_TABLE_ROWS = 12  # 11 KPIs + header
KPI_TABLE_COLS = 5  # KPI, Week 1, Week 2, Week 3, Week 4
CAPACITY_TABLE_ROWS = 4  # Header + 3 rows
CAPACITY_TABLE_COLS = 5

# Column width constants (inches)
KPI_COL_WIDTH = 2.5
WEEK_COL_WIDTH = 1.5

# Alternating row background color
ALT_ROW_COLOR = RGBColor(242, 242, 242)
HEADER_TEXT_COLOR = RGBColor(255, 255, 255)


def _format_minutes_value(minutes: float) -> str:
    """Format minutes value, converting to hours if >= 60.

    Args:
        minutes: Minutes value

    Returns:
        Formatted string (e.g., "2 Hour 30 Minutes" or "45 Minutes")
    """
    if minutes >= 60:
        hours = int(minutes // 60)
        mins = int(minutes % 60)
        return f"{hours} Hour {mins} Minutes"
    return f"{int(minutes)} Minutes"


def _format_kpi_value(
    value: Any, unit: str, include_change: bool = False, pct_change: str = ""
) -> str:
    """Format KPI value based on unit type.

    Args:
        value: KPI value
        unit: Unit type ("$", "Hours", "%", "Minutes", or empty)
        include_change: Whether to include percentage change
        pct_change: Percentage change string (if include_change is True)

    Returns:
        Formatted value string
    """
    if value is None:
        return ""
    change_suffix = f" {pct_change}" if include_change and pct_change else ""
    if unit == "$":
        return f"${int(value)}*{change_suffix}"
    elif unit == "Hours":
        return f"{value} {unit}{change_suffix}"
    elif unit == "%":
        return f"{value}{unit}{change_suffix}"
    elif unit == "Minutes":
        formatted = _format_minutes_value(value)
        return f"{formatted}{change_suffix}"
    else:
        return f"{int(value)} {unit}{change_suffix}" if value else ""


def _setup_kpi_table_title(slide: Any, ppt: PPTGenerator, equipment_code: str) -> None:
    """Setup table title textbox on the given slide.

    Args:
        slide: PowerPoint slide object
        ppt: PPTGenerator instance (used for color constants)
        equipment_code: Equipment identifier for the title
    """
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = f"Tooling Performance: {equipment_code}"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = ppt.PRIMARY_COLOR


def _setup_kpi_table_structure(slide: Any) -> Any:
    """Create KPI table shape and set column widths.

    Args:
        slide: PowerPoint slide object

    Returns:
        Table object with configured column widths
    """
    table = slide.shapes.add_table(
        KPI_TABLE_ROWS,
        KPI_TABLE_COLS,
        Inches(0.5),
        Inches(1.2),
        Inches(9),
        Inches(5.5),
    ).table

    # Set column widths
    table.columns[0].width = Inches(KPI_COL_WIDTH)
    for col_idx in range(1, KPI_TABLE_COLS):
        table.columns[col_idx].width = Inches(WEEK_COL_WIDTH)

    return table


def _setup_kpi_table_header(table: Any, ppt: PPTGenerator) -> None:
    """Setup table header row with styled column labels.

    Args:
        table: PowerPoint table object
        ppt: PPTGenerator instance (used for color constants)
    """
    headers = ["KPI", LABEL_WEEK_1, LABEL_WEEK_2, "Week 3", "Week 4"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ppt.PRIMARY_COLOR
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = HEADER_TEXT_COLOR
        paragraph.alignment = PP_ALIGN.CENTER


def _format_kpi_row_cells(table: Any, row_idx: int) -> None:
    """Format cells in a KPI row (alignment, font size, alternating colors).

    Args:
        table: PowerPoint table object
        row_idx: Row index to format
    """
    # Format value cells
    for col in range(1, KPI_TABLE_COLS):
        cell = table.cell(row_idx, col)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER

    # Alternate row colors
    if row_idx % 2 == 0:
        for col in range(KPI_TABLE_COLS):
            cell = table.cell(row_idx, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ALT_ROW_COLOR


def _add_kpi_table_footer(slide: Any, ppt: PPTGenerator) -> None:
    """Add footer note to KPI table slide explaining cost assumptions.

    Args:
        slide: PowerPoint slide object
        ppt: PPTGenerator instance (used for color constants)
    """
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    note_frame = note_box.text_frame
    note_frame.text = "*Based on a machine rate cost of $170/h and labor costs of $10/h"
    note_p = note_frame.paragraphs[0]
    note_p.font.size = Pt(8)
    note_p.font.italic = True
    note_p.font.color.rgb = ppt.LIGHT_GRAY


def _build_kpi_row_definitions() -> list:
    """Return the list of KPI row definitions for the comparison table.

    Returns:
        List of tuples: (display_name, kpi_key, unit, show_change)
    """
    return [
        ("Total Incurred Loss in Machine Hours:", "loss_machine_hours", "Hours", True),
        ("Total Incurred Costs (Machine Rate + Labor)*", "total_costs", "$", True),
        ("Parts Opportunity Lost", "parts_opportunity_lost", "", True),
        ("Run Rate Efficiency", "runrate_efficiency", "%", True),
        ("Run Rate MTBF", "runrate_mtbf_minutes", "Minutes", True),
        ("Run Rate MTTR", "runrate_mttr_minutes", "Minutes", True),
        ("Capacity Risk (OEE Score)", "capacity_oee_score", "%", False),
        ("Optimal Output (100% OEE)", "optimal_output", "parts", True),
        ("Actual Output", "actual_output", "parts", True),
        ("Availability (Gain/Loss)", "availability_gain_loss", "parts", True),
        ("Efficiency (Gain/Loss)", "efficiency_gain_loss", "parts", True),
    ]


def _populate_kpi_row(
    table: Any,
    row_idx: int,
    kpi_name: str,
    kpi_key: str,
    unit: str,
    week1_kpis: Dict[str, Any],
    week2_kpis: Dict[str, Any],
    ppt: PPTGenerator,
    calculate_percentage_change_fn: Any,
) -> None:
    """Populate a single KPI row with name, values, and formatting.

    Args:
        table: PowerPoint table object
        row_idx: Row index in the table
        kpi_name: Display name for the KPI
        kpi_key: Key to look up in KPI dictionaries
        unit: Unit string for formatting
        week1_kpis: Week 1 KPIs dictionary
        week2_kpis: Week 2 KPIs dictionary
        ppt: PPTGenerator instance (used for color constants)
        calculate_percentage_change_fn: Function to calculate percentage change
    """
    # KPI name
    cell = table.cell(row_idx, 0)
    cell.text = kpi_name
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = ppt.TEXT_COLOR

    # Week 1 value
    week1_val = week1_kpis.get(kpi_key)
    table.cell(row_idx, 1).text = _format_kpi_value(week1_val, unit)

    # Week 2 value with percentage change
    week2_val = week2_kpis.get(kpi_key)
    if week2_val is not None and week1_val is not None:
        pct_change = calculate_percentage_change_fn(week2_val, week1_val)
        table.cell(row_idx, 2).text = _format_kpi_value(
            week2_val, unit, include_change=True, pct_change=pct_change
        )
    else:
        table.cell(row_idx, 2).text = ""

    # Week 3 and 4 empty for now
    table.cell(row_idx, 3).text = ""
    table.cell(row_idx, 4).text = ""

    # Format cells
    _format_kpi_row_cells(table, row_idx)


def add_kpi_comparison_table(
    ppt: PPTGenerator,
    equipment_code: str,
    week1_kpis: Dict[str, Any],
    week2_kpis: Dict[str, Any],
    calculate_percentage_change_fn: Any,
) -> None:
    """Add KPI comparison table slide to the presentation.

    Args:
        ppt: PPTGenerator instance
        equipment_code: Equipment identifier
        week1_kpis: Week 1 KPIs dictionary
        week2_kpis: Week 2 KPIs dictionary
        calculate_percentage_change_fn: Function to calculate percentage change
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    # Setup title
    _setup_kpi_table_title(slide, ppt, equipment_code)

    # Create table structure
    table = _setup_kpi_table_structure(slide)

    # Setup header
    _setup_kpi_table_header(table, ppt)

    # Populate KPI rows
    kpi_rows = _build_kpi_row_definitions()
    for row_idx, (kpi_name, kpi_key, unit, _show_change) in enumerate(
        kpi_rows, start=1
    ):
        _populate_kpi_row(
            table,
            row_idx,
            kpi_name,
            kpi_key,
            unit,
            week1_kpis,
            week2_kpis,
            ppt,
            calculate_percentage_change_fn,
        )

    # Footer note
    _add_kpi_table_footer(slide, ppt)


def _populate_capacity_row(
    table: Any,
    row_idx: int,
    label: str,
    key: str,
    unit: str,
    week1_kpis: Dict[str, Any],
    week2_kpis: Dict[str, Any],
    calculate_percentage_change_fn: Any,
) -> None:
    """Populate a single capacity detail row with values, change, and trend.

    Args:
        table: PowerPoint table object
        row_idx: Row index in the table
        label: Display label for the metric
        key: Key to look up in KPI dictionaries
        unit: Unit string for formatting
        week1_kpis: Week 1 KPIs dictionary
        week2_kpis: Week 2 KPIs dictionary
        calculate_percentage_change_fn: Function to calculate percentage change
    """
    week1_val = int(week1_kpis.get(key, 0))
    week2_val = int(week2_kpis.get(key, 0))
    change = week2_val - week1_val
    pct_change = (
        calculate_percentage_change_fn(week2_val, week1_val)
        if week1_val != 0
        else "{N/A}"
    )

    table.cell(row_idx, 0).text = label
    table.cell(row_idx, 1).text = f"{week1_val} {unit}"
    table.cell(row_idx, 2).text = f"{week2_val} {unit}"
    table.cell(row_idx, 3).text = f"{change:+d} {unit} {pct_change}"

    # Trend indicator
    if change > 0:
        table.cell(row_idx, 4).text = "Improving"
    elif change < 0:
        table.cell(row_idx, 4).text = "Declining"
    else:
        table.cell(row_idx, 4).text = "Stable"

    # Format cells
    for col in range(CAPACITY_TABLE_COLS):
        cell = table.cell(row_idx, col)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT


def add_capacity_details_slide(
    ppt: PPTGenerator,
    equipment_code: str,
    week1_kpis: Dict[str, Any],
    week2_kpis: Dict[str, Any],
    calculate_percentage_change_fn: Any,
) -> None:
    """Add detailed capacity analysis slide to the presentation.

    Args:
        ppt: PPTGenerator instance
        equipment_code: Equipment identifier
        week1_kpis: Week 1 KPIs dictionary
        week2_kpis: Week 2 KPIs dictionary
        calculate_percentage_change_fn: Function to calculate percentage change
    """
    slide = ppt.prs.slides.add_slide(ppt.prs.slide_layouts[6])

    # Title
    ppt._add_slide_title(slide, f"Capacity Analysis Details: {equipment_code}")

    # Create detailed table
    table = slide.shapes.add_table(
        CAPACITY_TABLE_ROWS,
        CAPACITY_TABLE_COLS,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(3),
    ).table

    # Headers
    headers = ["Metric", LABEL_WEEK_1, LABEL_WEEK_2, "Change", "Trend"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ppt.PRIMARY_COLOR
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = HEADER_TEXT_COLOR
        paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    details = [
        ("Optimal Output (100% OEE)", "optimal_output", "parts"),
        ("Actual Output", "actual_output", "parts"),
        ("Output Gap", "output_gap", "parts"),
    ]

    for row_idx, (label, key, unit) in enumerate(details, start=1):
        _populate_capacity_row(
            table,
            row_idx,
            label,
            key,
            unit,
            week1_kpis,
            week2_kpis,
            calculate_percentage_change_fn,
        )
