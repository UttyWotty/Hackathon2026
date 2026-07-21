"""
Shared PowerPoint Generation Utilities
========================================

Provides common functions for creating professional PowerPoint presentations:
- Slide templates (title, content, chart, summary)
- Branding/theme management
- Chart insertion from files
- Table formatting
- Consistent styling

Author: Utku Gulbardak
Date: 2025-11-28
"""

import os
from datetime import datetime
from typing import Any, Dict, List, Optional

import pandas as pd  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.util import Inches, Pt  # type: ignore


class PPTGenerator:
    """
    Base class for PowerPoint generation with common utilities.

    Provides standardized slide templates and formatting for consistent,
    professional presentations across all analysis types.
    """

    # Brand colors (can be customized)
    PRIMARY_COLOR = RGBColor(0, 112, 192)  # Blue
    SECONDARY_COLOR = RGBColor(68, 114, 196)  # Light Blue
    SUCCESS_COLOR = RGBColor(0, 176, 80)  # Green
    WARNING_COLOR = RGBColor(255, 192, 0)  # Yellow
    DANGER_COLOR = RGBColor(255, 0, 0)  # Red
    TEXT_COLOR = RGBColor(64, 64, 64)  # Dark Gray
    LIGHT_GRAY = RGBColor(217, 217, 217)  # Light Gray

    def __init__(self):
        """Initialize PowerPoint generator with blank presentation."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)  # Standard width
        self.prs.slide_height = Inches(7.5)  # Standard height

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        metadata: Optional[Dict[str, str]] = None,
    ) -> None:
        """
        Create title slide with branding.

        Args:
            title: Main title text
            subtitle: Optional subtitle text
            metadata: Optional metadata dict (equipment, date range, etc.)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add title (centered, large)
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = self.PRIMARY_COLOR

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.5), Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.alignment = PP_ALIGN.CENTER
            subtitle_p.font.size = Pt(24)
            subtitle_p.font.color.rgb = self.SECONDARY_COLOR

        # Add metadata at bottom
        if metadata:
            y_pos = 5.5
            for key, value in metadata.items():
                meta_box = slide.shapes.add_textbox(
                    Inches(2), Inches(y_pos), Inches(6), Inches(0.4)
                )
                meta_frame = meta_box.text_frame
                meta_frame.text = f"{key}: {value}"
                meta_p = meta_frame.paragraphs[0]
                meta_p.alignment = PP_ALIGN.CENTER
                meta_p.font.size = Pt(14)
                meta_p.font.color.rgb = self.TEXT_COLOR
                y_pos += 0.4

        # Add timestamp footer
        timestamp_box = slide.shapes.add_textbox(
            Inches(1), Inches(7), Inches(8), Inches(0.3)
        )
        timestamp_frame = timestamp_box.text_frame
        timestamp_frame.text = (
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        )
        timestamp_p = timestamp_frame.paragraphs[0]
        timestamp_p.alignment = PP_ALIGN.CENTER
        timestamp_p.font.size = Pt(10)
        timestamp_p.font.color.rgb = self.LIGHT_GRAY

    def add_summary_slide(
        self, title: str, metrics: Dict[str, Any], layout: str = "grid"
    ) -> None:
        """
        Create executive summary slide with key metrics.

        Args:
            title: Slide title
            metrics: Dictionary of metric name -> value pairs
            layout: "grid" or "list" layout style
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add title
        self._add_slide_title(slide, title)

        if layout == "grid":
            self._add_metrics_grid(slide, metrics)
        else:
            self._add_metrics_list(slide, metrics)

    def _add_slide_title(self, slide, title: str) -> None:
        """Add consistent title to slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = self.PRIMARY_COLOR

    def _add_metrics_grid(self, slide, metrics: Dict[str, Any]) -> None:
        """Add metrics in grid layout (2x2 or 2x3)."""
        metrics_list = list(metrics.items())
        cols = 2

        box_width = 4
        box_height = 1.5
        x_start = 1
        y_start = 1.5
        x_gap = 0.5
        y_gap = 0.5

        for idx, (key, value) in enumerate(metrics_list):
            row = idx // cols
            col = idx % cols

            x = x_start + col * (box_width + x_gap)
            y = y_start + row * (box_height + y_gap)

            # Metric box background
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(x),
                Inches(y),
                Inches(box_width),
                Inches(box_height),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
            shape.line.color.rgb = self.PRIMARY_COLOR

            # Metric value (large)
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 0.2), Inches(box_width), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = str(value)
            value_p = value_frame.paragraphs[0]
            value_p.alignment = PP_ALIGN.CENTER
            value_p.font.size = Pt(36)
            value_p.font.bold = True
            value_p.font.color.rgb = self.PRIMARY_COLOR

            # Metric label (small)
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 0.9), Inches(box_width), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = key
            label_p = label_frame.paragraphs[0]
            label_p.alignment = PP_ALIGN.CENTER
            label_p.font.size = Pt(14)
            label_p.font.color.rgb = self.TEXT_COLOR

    def _add_metrics_list(self, slide, metrics: Dict[str, Any]) -> None:
        """Add metrics in list layout."""
        y_pos = 1.5
        for key, value in metrics.items():
            # Label
            label_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_pos), Inches(4), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = key
            label_p = label_frame.paragraphs[0]
            label_p.font.size = Pt(18)
            label_p.font.color.rgb = self.TEXT_COLOR

            # Value
            value_box = slide.shapes.add_textbox(
                Inches(5), Inches(y_pos), Inches(3), Inches(0.4)
            )
            value_frame = value_box.text_frame
            value_frame.text = str(value)
            value_p = value_frame.paragraphs[0]
            value_p.alignment = PP_ALIGN.RIGHT
            value_p.font.size = Pt(18)
            value_p.font.bold = True
            value_p.font.color.rgb = self.PRIMARY_COLOR

            y_pos += 0.6

    def add_chart_slide(
        self, title: str, image_path: str, caption: Optional[str] = None
    ) -> None:
        """
        Insert chart/image slide.

        Args:
            title: Slide title
            image_path: Path to chart image file
            caption: Optional caption text
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        self._add_slide_title(slide, title)

        # Add image if file exists
        if os.path.exists(image_path):
            # Center the image
            img_width = Inches(8)
            img_height = Inches(5)
            img_left = Inches(1)
            img_top = Inches(1.2)

            slide.shapes.add_picture(
                image_path, img_left, img_top, width=img_width, height=img_height
            )
        else:
            # Placeholder if image not found
            placeholder_box = slide.shapes.add_textbox(
                Inches(3), Inches(3.5), Inches(4), Inches(0.5)
            )
            placeholder_frame = placeholder_box.text_frame
            placeholder_frame.text = (
                f"Chart not available: {os.path.basename(image_path)}"
            )
            placeholder_p = placeholder_frame.paragraphs[0]
            placeholder_p.alignment = PP_ALIGN.CENTER
            placeholder_p.font.size = Pt(14)
            placeholder_p.font.color.rgb = self.LIGHT_GRAY

        # Add caption if provided
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1), Inches(6.5), Inches(8), Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_p = caption_frame.paragraphs[0]
            caption_p.alignment = PP_ALIGN.CENTER
            caption_p.font.size = Pt(12)
            caption_p.font.italic = True
            caption_p.font.color.rgb = self.TEXT_COLOR

    def add_table_slide(
        self, title: str, data: pd.DataFrame, max_rows: int = 10
    ) -> None:
        """
        Create slide with formatted table.

        Args:
            title: Slide title
            data: DataFrame to display
            max_rows: Maximum rows to show (will truncate if needed)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        self._add_slide_title(slide, title)

        # Truncate data if needed
        display_data = data.head(max_rows) if len(data) > max_rows else data

        # Calculate table dimensions
        rows = len(display_data) + 1  # +1 for header
        cols = len(display_data.columns)

        # Add table
        table_left = Inches(0.5)
        table_top = Inches(1.2)
        table_width = Inches(9)
        table_height = Inches(5.5)

        table = slide.shapes.add_table(
            rows, cols, table_left, table_top, table_width, table_height
        ).table

        # Set column headers
        for col_idx, col_name in enumerate(display_data.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.PRIMARY_COLOR
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Fill data rows
        for row_idx, row in enumerate(display_data.itertuples(index=False), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = self.TEXT_COLOR

                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

        # Add truncation note if data was truncated
        if len(data) > max_rows:
            note_box = slide.shapes.add_textbox(
                Inches(1), Inches(7), Inches(8), Inches(0.3)
            )
            note_frame = note_box.text_frame
            note_frame.text = f"Showing {max_rows} of {len(data)} rows"
            note_p = note_frame.paragraphs[0]
            note_p.alignment = PP_ALIGN.CENTER
            note_p.font.size = Pt(10)
            note_p.font.italic = True
            note_p.font.color.rgb = self.LIGHT_GRAY

    def add_text_slide(
        self, title: str, content: List[str], bullet: bool = True
    ) -> None:
        """
        Add text content slide with optional bullet points.

        Args:
            title: Slide title
            content: List of text items
            bullet: Whether to use bullet points
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        self._add_slide_title(slide, title)

        # Add content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for idx, item in enumerate(content):
            if idx > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = item
            p.level = 1 if bullet else 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)

    def save(self, output_path: str) -> str:
        """
        Save presentation to file.

        Args:
            output_path: Path to save PPT file

        Returns:
            Path to saved file
        """
        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Save presentation
        self.prs.save(output_path)

        return output_path


def format_metric_value(value: Any) -> str:
    """
    Format metric value for display in presentation.

    Args:
        value: Raw metric value

    Returns:
        Formatted string
    """
    if isinstance(value, float):
        if value >= 1000:
            return f"{value:,.0f}"
        elif value >= 1:
            return f"{value:.2f}"
        else:
            return f"{value:.4f}"
    elif isinstance(value, int):
        return f"{value:,}"
    else:
        return str(value)
